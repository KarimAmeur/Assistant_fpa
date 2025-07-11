import streamlit as st
import os
import tempfile
import time
import logging
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.vectorstores import Chroma
import pandas as pd
import PyPDF2
import docx
import torch
from langchain_community.embeddings import HuggingFaceEmbeddings

# Import pour PowerPoint
try:
    from pptx import Presentation
except ImportError:
    # Tenter d'installer pptx si non disponible
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

# Configuration du logging
logging.basicConfig(level=logging.INFO,
                   format='%(asctime)s - %(levelname)s - %(message)s')

def extract_text_from_pdf(file_path: str) -> str:
    """Extraire le texte d'un fichier PDF."""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            full_text = ""
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + "\n\n"
            return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du PDF {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extraire le texte d'un fichier Word."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extraire le texte d'un fichier PowerPoint."""
    try:
        prs = Presentation(file_path)
        full_text = ""
        
        # Extraire le texte de chaque diapositive
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extraire le titre de la diapositive
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text
                if title_text.strip():
                    slide_text.append(f"Titre: {title_text}")
            
            # Extraire le texte de toutes les formes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                elif hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text)
            
            # Ajouter le texte de la diapositive au texte global
            if slide_text:
                full_text += f"Diapositive {i}:\n"
                full_text += "\n".join(slide_text) + "\n\n"
        
        return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du PowerPoint {file_path}: {e}")
        return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extraire le texte d'un fichier Excel."""
    try:
        # Lire tous les onglets
        xls = pd.ExcelFile(file_path)
        full_text = ""

        for sheet_name in xls.sheet_names:
            # Lire la feuille
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Convertir toutes les colonnes en chaîne et concaténer
            sheet_text = "\n".join([
                " ".join(str(val) for val in row if pd.notna(val))
                for row in df.values
            ])

            full_text += f"Feuille {sheet_name}:\n{sheet_text}\n\n"

        return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction de l'XLSX {file_path}: {e}")
        return ""

def user_rag_page():
    """Page permettant à l'utilisateur d'ajouter ses propres documents au RAG."""
    
    # Bannière avec logo et titre
    st.markdown("""
    <div class="banner">
        <h1>📄 Importation de Documents Personnels</h1>
        <p>Ajoutez vos propres documents pour enrichir la base de connaissances de l'assistant</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Vérification de la disponibilité des ressources nécessaires
    if 'RAG_user' not in st.session_state:
        st.session_state.RAG_user = None
    
    # Colonnes pour l'interface
    left_col, right_col = st.columns([3, 2])
    
    with left_col:
        st.markdown("""
        <div class="scenario-card">
            <h3>🔍 Importation de Documents</h3>
            <p>Téléchargez vos documents pour enrichir la base de connaissances personnalisée.</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Widget de téléchargement de fichiers
        uploaded_files = st.file_uploader(
            "Sélectionnez un ou plusieurs fichiers (PDF, DOCX, PPTX, XLSX)",
            type=["pdf", "docx", "pptx", "ppt", "xlsx", "xls"],
            accept_multiple_files=True,
            help="Les fichiers seront traités puis ajoutés à votre RAG personnel"
        )
        
        # Paramètres fixes pour les chunks
        chunk_size = 256
        chunk_overlap_percent = 10
        overlap_chars = int(chunk_size * (chunk_overlap_percent / 100))
        
        st.info(f"Les documents seront automatiquement découpés en chunks de {chunk_size} caractères, avec un chevauchement de {overlap_chars} caractères.")
        
        # Bouton de traitement
        process_button = st.button(
            "📊 Traiter et vectoriser les documents",
            type="primary",
            use_container_width=True,
            disabled=not uploaded_files
        )
        
        # Traitement des fichiers
        if process_button and uploaded_files:
            with st.status("Traitement des documents en cours...", expanded=True) as status:
                # Création d'un dossier temporaire pour enregistrer les fichiers
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Sauvegarde des fichiers dans le dossier temporaire
                    file_paths = []
                    for uploaded_file in uploaded_files:
                        file_path = os.path.join(temp_dir, uploaded_file.name)
                        with open(file_path, "wb") as f:
                            f.write(uploaded_file.getbuffer())
                        file_paths.append(file_path)
                        st.write(f"Fichier sauvegardé : {uploaded_file.name}")
                    
                    # Extraction du texte des fichiers
                    documents = []
                    
                    for file_path in file_paths:
                        st.write(f"Traitement de {os.path.basename(file_path)}...")
                        
                        # Déterminer le type de fichier et extraire le texte
                        file_extension = os.path.splitext(file_path)[1].lower()
                        
                        if file_extension == '.pdf':
                            text = extract_text_from_pdf(file_path)
                        elif file_extension == '.docx':
                            text = extract_text_from_docx(file_path)
                        elif file_extension in ['.pptx', '.ppt']:
                            text = extract_text_from_pptx(file_path)
                        elif file_extension in ['.xlsx', '.xls']:
                            text = extract_text_from_xlsx(file_path)
                        else:
                            st.warning(f"Type de fichier non supporté : {file_extension}")
                            continue
                        
                        if not text.strip():
                            st.warning(f"Aucun texte extrait de {os.path.basename(file_path)}")
                            continue
                        
                        st.write(f"✅ Texte extrait ({len(text)} caractères)")
                        
                        # Découpage du texte en chunks avec paramètres fixes
                        text_splitter = RecursiveCharacterTextSplitter(
                            chunk_size=chunk_size,
                            chunk_overlap=overlap_chars,
                            length_function=len,
                            separators=["\n\n", "\n", ".", "!", "?", ";", ":", " ", ""]
                        )
                        
                        chunks = text_splitter.split_text(text)
                        st.write(f"🔪 Texte découpé en {len(chunks)} chunks")
                        
                        # Création des documents Langchain
                        filename = os.path.basename(file_path)
                        for i, chunk in enumerate(chunks):
                            if len(chunk.strip()) < 30:  # Éviter les chunks trop courts
                                continue
                            
                            documents.append(Document(
                                page_content=chunk,
                                metadata={
                                    "source": filename,
                                    "filename": filename,
                                    "chunk": i,
                                    "type": file_extension[1:],  # Type sans le point initial
                                    "size": len(chunk)
                                }
                            ))
                    
                    # Création ou mise à jour de la base vectorielle
                    if documents:
                        st.write(f"🧠 Vectorisation de {len(documents)} documents...")
                        
                        try:
                            # Récupérer le modèle d'embedding
                            embeddings = get_embedding_model()
                            
                            # Déterminer s'il faut créer une nouvelle base ou mettre à jour l'existante
                            if st.session_state.RAG_user is None:
                                # Création d'une nouvelle base vectorielle
                                persist_directory = "chroma_db_user"
                                st.session_state.RAG_user = Chroma.from_documents(
                                    documents=documents,
                                    embedding=embeddings,
                                    persist_directory=persist_directory
                                )
                                st.session_state.RAG_user.persist()
                                st.write(f"🎉 Nouvelle base vectorielle créée avec {len(documents)} documents")
                            else:
                                # Mise à jour de la base existante
                                st.session_state.RAG_user.add_documents(documents=documents)
                                st.session_state.RAG_user.persist()
                                st.write(f"🔄 Base vectorielle mise à jour avec {len(documents)} nouveaux documents")
                            
                            status.update(label="Traitement terminé avec succès!", state="complete", expanded=False)
                            
                            # Indiquer le nombre total de documents dans la base
                            total_docs = len(st.session_state.RAG_user.get())
                            st.success(f"Votre RAG personnel contient maintenant {total_docs} chunks de documents")
                            st.rerun()
                        except Exception as e:
                            status.update(label=f"Erreur lors de la vectorisation: {str(e)}", state="error")
                            st.error(f"Erreur de vectorisation: {str(e)}")
                    else:
                        status.update(label="Aucun document valide à vectoriser", state="error")
                        st.warning("Aucun document valide n'a pu être traité")
    
    
    with right_col:
        st.markdown("""
        <div class="scenario-card">
            <h3>ℹ️ Informations sur votre RAG personnel</h3>
        </div>
        """, unsafe_allow_html=True)
        
        # Affichage des informations sur le RAG personnel
        if st.session_state.RAG_user is not None:
            try:
                # Récupérer les informations sur la base vectorielle
                docs_dict = st.session_state.RAG_user.get()
                total_docs = len(docs_dict['documents'])
            
                # Récupérer les métadonnées des documents
                sources = {}
                file_types = {}
                
                for metadata in docs_dict['metadatas']:
                    filename = metadata.get('filename', 'Inconnu')
                    filetype = metadata.get('type', 'Inconnu')
                    
                    # Compter par nom de fichier
                    if filename in sources:
                        sources[filename] += 1
                    else:
                        sources[filename] = 1
                    
                    # Compter par type de fichier
                    if filetype in file_types:
                        file_types[filetype] += 1
                    else:
                        file_types[filetype] = 1
                
                # Le reste du code reste similaire...
            except Exception as e:
                st.error(f"Erreur lors de la récupération des informations sur le RAG: {str(e)}")
                
                st.markdown(f"""
                <div class="info-box">
                    <p><strong>📊 Statistiques de votre RAG personnel:</strong></p>
                    <ul>
                        <li><strong>Nombre total de chunks:</strong> {total_docs}</li>
                        <li><strong>Nombre de fichiers sources:</strong> {len(sources)}</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)
                
                # Afficher la répartition par type de fichier
                st.markdown("<p><strong>📊 Répartition par type de fichier:</strong></p>", unsafe_allow_html=True)
                for filetype, count in file_types.items():
                    file_emoji = {
                        'pdf': '📄', 
                        'docx': '📝', 
                        'pptx': '📊',
                        'ppt': '📊',
                        'xlsx': '📈',
                        'xls': '📈'
                    }.get(filetype, '📁')
                    
                    st.markdown(f"- {file_emoji} **{filetype.upper()}**: {count} chunks")
                
                # Afficher la liste des fichiers
                st.markdown("<p><strong>📁 Fichiers vectorisés:</strong></p>", unsafe_allow_html=True)
                for filename, count in sources.items():
                    st.markdown(f"- **{filename}**: {count} chunks")
                
                # Option pour tester le RAG
                st.markdown("""
                <div class="scenario-card" style="margin-top: 20px;">
                    <h3>🔍 Tester votre RAG personnel</h3>
                </div>
                """, unsafe_allow_html=True)
                
                test_query = st.text_input(
                    "Saisissez une requête de test",
                    placeholder="Exemple: Quels sont les points clés abordés dans mes documents?",
                    help="Cette requête sera utilisée pour rechercher dans vos documents"
                )
                
                if st.button("🔎 Rechercher", use_container_width=True) and test_query:
                    # Effectuer la recherche
                    with st.spinner("Recherche en cours..."):
                        # Effectuer la recherche similaire
                        results = st.session_state.RAG_user.similarity_search_with_score(
                        query=test_query,
                        k=3  # Nombre de résultats à afficher
                        )
                        
                        if results:
                            st.markdown("<p><strong>📑 Résultats de la recherche:</strong></p>", unsafe_allow_html=True)
                            for i, (doc, score) in enumerate(results, 1):
                                file_type = doc.metadata.get('type', 'inconnu')
                                file_emoji = {
                                    'pdf': '📄', 
                                    'docx': '📝', 
                                    'pptx': '📊',
                                    'ppt': '📊',
                                    'xlsx': '📈',
                                    'xls': '📈'
                                }.get(file_type, '📁')
                                
                                with st.expander(f"{file_emoji} Résultat {i} - Score: {score:.4f} - Source: {doc.metadata.get('filename', 'Inconnu')}"):
                                    st.markdown(f"**Extrait du document:**\n\n{doc.page_content}")
                        else:
                            st.info("Aucun résultat correspondant à votre requête")
            except Exception as e:
                st.error(f"Erreur lors de la récupération des informations sur le RAG: {str(e)}")
        else:
            st.info("Aucun document n'a encore été importé dans votre RAG personnel")
            
            st.markdown("""
            <div class="info-box" style="margin-top: 20px;">
                <p><strong>ℹ️ Comment utiliser cette fonctionnalité:</strong></p>
                <ol>
                    <li>Téléchargez vos documents (PDF, Word, PowerPoint, Excel)</li>
                    <li>Cliquez sur "Traiter et vectoriser les documents"</li>
                    <li>Une fois vos documents vectorisés, vous pourrez tester la recherche</li>
                </ol>
                <p>Vos documents seront stockés uniquement dans votre session Streamlit actuelle et ne seront pas conservés après la fermeture de l'application.</p>
            </div>
            """, unsafe_allow_html=True)
        
        # Ajouter un bouton pour effacer le RAG personnel
        if st.session_state.RAG_user is not None:
            if st.button("🗑️ Effacer mon RAG personnel", type="secondary", use_container_width=True):
                # Confirmation
                confirm = st.checkbox("Confirmer la suppression de tous mes documents")
                if confirm:
                    st.session_state.RAG_user = None
                    st.success("Votre RAG personnel a été effacé")
                    st.experimental_rerun()